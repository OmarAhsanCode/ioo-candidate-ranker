import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5) # 16:9 aspect ratio

    # Colors
    NAVY = RGBColor(10, 20, 38)
    WHITE = RGBColor(255, 255, 255)
    CYAN = RGBColor(0, 198, 255)
    LIGHT_GRAY = RGBColor(180, 190, 210)
    DARK_GRAY = RGBColor(20, 35, 60)
    BORDER_COLOR = RGBColor(0, 198, 255)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY

    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = CYAN
        return title_box

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1)
    
    # Title & Subtitle text frame
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = "Making Hiring Smarter"
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(64)
    p_title.font.bold = True
    p_title.font.color.rgb = CYAN
    p_title.space_after = Pt(20)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "An AI-Powered Candidate Ranking System for Senior AI Engineer @ Redrob"
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = WHITE
    
    # Footer info
    footer_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.33), Inches(1.0))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_bottom = tf_f.margin_right = 0
    p_f1 = tf_f.paragraphs[0]
    p_f1.text = "Team IOO | Iqra · Omar · Ojas"
    p_f1.font.name = 'Segoe UI'
    p_f1.font.size = Pt(16)
    p_f1.font.bold = True
    p_f1.font.color.rgb = CYAN
    
    p_f2 = tf_f.add_paragraph()
    p_f2.text = "India Runs at Active Skills — Data & AI Challenge 2026"
    p_f2.font.name = 'Segoe UI'
    p_f2.font.size = Pt(14)
    p_f2.font.color.rgb = LIGHT_GRAY

    # Slide 2: The Problem
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2)
    add_title(slide2, "Keyword Filters Fail Great Candidates")
    
    content_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
    
    points = [
        ("Massive Scale Bottlenecks", "Recruiters screen 100,000+ profiles manually — a volume impossible to filter thoroughly at scale without losing diamonds in the rough."),
        ("Buzzword Optimization Abuse", "Traditional keyword matching rewards candidates who deliberately stuff resumes with trending buzzwords, rather than demonstrating genuine, verified talent."),
        ("Missing the Best Candidates", "The most capable engineers often don't write or structure their profiles specifically for ATS parsers, causing them to get auto-rejected.")
    ]
    
    for title, desc in points:
        p_t = tf2.add_paragraph() if tf2.text else tf2.paragraphs[0]
        p_t.text = f"▪ {title}"
        p_t.font.name = 'Segoe UI'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN
        p_t.space_before = Pt(20)
        p_t.space_after = Pt(5)
        
        p_d = tf2.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Segoe UI'
        p_d.font.size = Pt(18)
        p_d.font.color.rgb = WHITE
        p_d.space_after = Pt(15)

    # Slide 3: Our Approach (Overview Diagram)
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3)
    add_title(slide3, "A Two-Stage Intelligent Pipeline")
    
    # Draw Left-to-Right Boxes
    box_width = Inches(2.2)
    box_height = Inches(1.2)
    y_pos = Inches(2.2)
    start_x = Inches(0.75)
    gap = Inches(0.8)
    
    steps = ["Job Description", "Feature Extraction", "Composite Scorer", "Ranked Shortlist"]
    
    for i, step in enumerate(steps):
        x_pos = start_x + i * (box_width + gap)
        # Create shape
        shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_GRAY
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(2)
        
        tf_s = shape.text_frame
        tf_s.word_wrap = True
        p_s = tf_s.paragraphs[0]
        p_s.text = step
        p_s.alignment = PP_ALIGN.CENTER
        p_s.font.name = 'Segoe UI'
        p_s.font.size = Pt(18)
        p_s.font.bold = True
        p_s.font.color.rgb = WHITE
        
        # Draw arrow to next box if not the last box
        if i < len(steps) - 1:
            arrow_x = x_pos + box_width
            arrow_w = gap
            arrow = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x + Inches(0.1), y_pos + Inches(0.4), arrow_w - Inches(0.2), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = CYAN
            arrow.line.color.rgb = CYAN
            
    # Bullet points below diagram
    desc_box = slide3.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.83), Inches(2.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_top = tf_desc.margin_bottom = tf_desc.margin_right = 0
    
    p1 = tf_desc.paragraphs[0]
    p1.text = "▪ Stage 1: Parse all 100K candidates and extract structured feature vectors"
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = CYAN
    p1.space_after = Pt(10)
    
    p1_sub = tf_desc.add_paragraph()
    p1_sub.text = "   Runs honeypot detection rules and disqualification layer filtering to immediately eliminate invalid/unsuited candidates."
    p1_sub.font.name = 'Segoe UI'
    p1_sub.font.size = Pt(16)
    p1_sub.font.color.rgb = WHITE
    p1_sub.space_after = Pt(20)
    
    p2 = tf_desc.add_paragraph()
    p2.text = "▪ Stage 2: Apply weighted scoring formula, sort, output top 100 with reasoning"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = CYAN
    p2.space_after = Pt(10)
    
    p2_sub = tf_desc.add_paragraph()
    p2_sub.text = "   Integrates technical capability with behavioral signals to generate an actionable rank shortlist sorted by exact criteria."
    p2_sub.font.name = 'Segoe UI'
    p2_sub.font.size = Pt(16)
    p2_sub.font.color.rgb = WHITE

    # Slide 4: What We Measure
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4)
    add_title(slide4, "Five Signal Layers — Weighted by Relevance")
    
    # Add Table
    rows = 6
    cols = 2
    left = Inches(0.75)
    top = Inches(1.8)
    width = Inches(11.83)
    height = Inches(3.5)
    
    table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(8.83)
    table.columns[1].width = Inches(3.0)
    
    # Table Content
    data = [
        ("Signal Layer", "Weight"),
        ("Core Skills Match (Vector DBs, RAG, Information Retrieval, Python)", "35%"),
        ("Career Trajectory & Stability (Product experience, role progression, tenure)", "30%"),
        ("Years of Experience (Linear scale up to target 8+ years)", "15%"),
        ("Location Fit (India-based + willingness to relocate)", "15%"),
        ("Education (CS/ML tier-based ranking as a small tiebreaker)", "5%")
    ]
    
    for row_idx, row in enumerate(data):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_GRAY if row_idx > 0 else NAVY
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Segoe UI'
            p.font.size = Pt(18) if row_idx > 0 else Pt(20)
            p.font.bold = True if row_idx == 0 or col_idx == 1 else False
            p.font.color.rgb = CYAN if row_idx == 0 or col_idx == 1 else WHITE
            
    # Bottom Note
    note_box = slide4.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(11.83), Inches(1.0))
    tf_n = note_box.text_frame
    tf_n.word_wrap = True
    tf_n.margin_left = tf_n.margin_top = tf_n.margin_bottom = tf_n.margin_right = 0
    p_n = tf_n.paragraphs[0]
    p_n.text = "Final Score = JD Match Score × Behavioral Availability Multiplier"
    p_n.font.name = 'Segoe UI'
    p_n.font.size = Pt(22)
    p_n.font.bold = True
    p_n.font.color.rgb = CYAN

    # Slide 5: Skills Scoring
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5)
    add_title(slide5, "Skills: Beyond Keyword Matching")
    
    # Left Column: Must Have
    left_box = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5.0))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_right = 0
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "MUST HAVE (High Weight)"
    p_lh.font.name = 'Segoe UI'
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = CYAN
    p_lh.space_after = Pt(15)
    
    skills_list = [
        "Vector DBs: Pinecone, Weaviate, Qdrant, FAISS, pgvector, OpenSearch",
        "Retrieval: BM25, Hybrid Search, Dense/Sparse Retrieval, Information Retrieval",
        "Evaluation: NDCG, MRR, A/B Testing, Search Relevance",
        "Production AI: RAG, Large Language Models (LLMs), Python"
    ]
    for s in skills_list:
        p_s = tf_l.add_paragraph()
        p_s.text = f"▪ {s}"
        p_s.font.name = 'Segoe UI'
        p_s.font.size = Pt(18)
        p_s.font.color.rgb = WHITE
        p_s.space_after = Pt(15)

    # Right Column: How we score
    right_box = slide5.shapes.add_textbox(Inches(6.83), Inches(1.8), Inches(5.75), Inches(5.0))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_right = 0
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "How We Score"
    p_rh.font.name = 'Segoe UI'
    p_rh.font.size = Pt(22)
    p_rh.font.bold = True
    p_rh.font.color.rgb = CYAN
    p_rh.space_after = Pt(15)
    
    scoring_points = [
        ("✓ Proficiency Scaling", "Scores scaled dynamically based on beginner → expert capability levels."),
        ("✓ Endorsement Counts", "Number of peer endorsements weight the candidate's core expertise."),
        ("✓ Assessment Score Bonus", "Scores of ≥70 on relevant technical assessments award a +0.15 bonus."),
        ("✗ Domain Disqualifications", "Primary expertise in Computer Vision, Robotics, or Speech is penalized/filtered to maintain focus on search engineering.")
    ]
    for title, desc in scoring_points:
        p_st = tf_r.add_paragraph()
        p_st.text = title
        p_st.font.name = 'Segoe UI'
        p_st.font.size = Pt(18)
        p_st.font.bold = True
        p_st.font.color.rgb = WHITE
        p_st.space_after = Pt(2)
        
        p_sd = tf_r.add_paragraph()
        p_sd.text = desc
        p_sd.font.name = 'Segoe UI'
        p_sd.font.size = Pt(14)
        p_sd.font.color.rgb = LIGHT_GRAY
        p_sd.space_after = Pt(10)

    # Slide 6: Career Trajectory Scoring
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6)
    add_title(slide6, "Career: We Read Between the Titles")
    
    # Grid Layout for Career Signal
    grid_data = [
        ("🏢 Product vs Services", "Candidates from pure IT services or consulting backgrounds are filtered/penalized to favor product-focused software engineering culture."),
        ("📈 Title Progression", "Rewarding candidates demonstrating natural career evolution: Software Engineer → Senior MLE → Lead AI Engineer."),
        ("⏱ Tenure Stability", "Consecutive short-term roles under 18 months trigger tenure penalties, ensuring long-term candidate commitment."),
        ("🎯 Role Relevance", "Targeted scoring boosts for current roles in Machine Learning, NLP, Information Retrieval, and Search Engine fields.")
    ]
    
    for idx, (title, desc) in enumerate(grid_data):
        row = idx // 2
        col = idx % 2
        col_x = Inches(0.75) if col == 0 else Inches(6.83)
        row_y = Inches(1.8) if row == 0 else Inches(4.3)
        
        box = slide6.shapes.add_textbox(col_x, row_y, Inches(5.75), Inches(2.2))
        tf_g = box.text_frame
        tf_g.word_wrap = True
        tf_g.margin_left = tf_g.margin_top = tf_g.margin_bottom = tf_g.margin_right = 0
        
        p_t = tf_g.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Segoe UI'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN
        p_t.space_after = Pt(8)
        
        p_d = tf_g.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Segoe UI'
        p_d.font.size = Pt(16)
        p_d.font.color.rgb = WHITE

    # Slide 7: Behavioral Availability
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7)
    add_title(slide7, "Great Candidate, Wrong Moment? We Factor That In.")
    
    # Scale box
    scale_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(11.83), Inches(1.5))
    scale_box.fill.solid()
    scale_box.fill.fore_color.rgb = DARK_GRAY
    scale_box.line.color.rgb = BORDER_COLOR
    scale_box.line.width = Pt(2)
    tf_scale = scale_box.text_frame
    tf_scale.word_wrap = True
    p_sc = tf_scale.paragraphs[0]
    p_sc.alignment = PP_ALIGN.CENTER
    p_sc.text = "Behavioral Availability Scale: 0.3x (Inactive / Hard Notice)  ◀──────────────────▶  1.0x (Highly Active / Immediate)"
    p_sc.font.name = 'Segoe UI'
    p_sc.font.size = Pt(20)
    p_sc.font.bold = True
    p_sc.font.color.rgb = CYAN
    
    # Metrics
    metrics_box = slide7.shapes.add_textbox(Inches(0.75), Inches(3.8), Inches(11.83), Inches(3.0))
    tf_m = metrics_box.text_frame
    tf_m.word_wrap = True
    tf_m.margin_left = tf_m.margin_top = tf_m.margin_bottom = tf_m.margin_right = 0
    
    p_mh = tf_m.paragraphs[0]
    p_mh.text = "Availability Signals Utilized:"
    p_mh.font.name = 'Segoe UI'
    p_mh.font.size = Pt(22)
    p_mh.font.bold = True
    p_mh.font.color.rgb = WHITE
    p_mh.space_after = Pt(15)
    
    sig_points = [
        ("• Days Since Last Active", "Decays from 1.0 (<30 days active) down to a base of 0.3 (>180 days inactive)."),
        ("• Notice Period Penalty", "Notice periods <= 30 days are rewarded (+0.04), while notice periods > 90 days are penalized (-0.06)."),
        ("• Engagement & Verification", "Recruiter response rates, interview completion rates, and verified contact details scale engagement scores.")
    ]
    
    for title, desc in sig_points:
        p_pt = tf_m.add_paragraph()
        p_pt.text = f"{title}: {desc}"
        p_pt.font.name = 'Segoe UI'
        p_pt.font.size = Pt(18)
        p_pt.font.color.rgb = LIGHT_GRAY
        p_pt.space_after = Pt(10)

    # Slide 8: Honeypot & Quality Filtering
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8)
    add_title(slide8, "We Catch What Filters Miss")
    
    # Honeypot rules
    rules_box = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf_rules = rules_box.text_frame
    tf_rules.word_wrap = True
    tf_rules.margin_left = tf_rules.margin_top = tf_rules.margin_bottom = tf_rules.margin_right = 0
    
    honeypots = [
        ("1. Timeline Impossibilities", "Disqualifies candidates listing employment start dates that occur before their plausible birth year (e.g. working at age 5)."),
        ("2. Skill Contradictions", "Filters profiles listing self-claimed 'Expert' proficiency levels on a technology with 0 months of actual usage history."),
        ("3. Experience Inflation", "Identifies candidates claiming a total years of experience that is mathematically greater than 2x their calculated career timeline."),
        ("4. Impossible Completeness", "Flags candidates claiming a 100% profile completeness metric despite leaving major mandatory profile fields entirely empty.")
    ]
    
    for title, desc in honeypots:
        p_ht = tf_rules.add_paragraph() if tf_rules.text else tf_rules.paragraphs[0]
        p_ht.text = title
        p_ht.font.name = 'Segoe UI'
        p_ht.font.size = Pt(20)
        p_ht.font.bold = True
        p_ht.font.color.rgb = CYAN
        p_ht.space_after = Pt(2)
        
        p_hd = tf_rules.add_paragraph()
        p_hd.text = desc
        p_hd.font.name = 'Segoe UI'
        p_hd.font.size = Pt(16)
        p_hd.font.color.rgb = WHITE
        p_hd.space_after = Pt(15)

    # Slide 9: Results (Top 10)
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9)
    add_title(slide9, "Our Top 10 — Clean, Defensible, Diverse")
    
    # Table layout
    rows = 11
    cols = 4
    t_left = Inches(0.75)
    t_top = Inches(1.5)
    t_width = Inches(11.83)
    t_height = Inches(5.2)
    
    table_shape = slide9.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
    table = table_shape.table
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(5.43)
    
    headers = ["Rank", "Title", "Location", "Key Strength"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Segoe UI'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
    top_10 = [
        ("1", "Applied ML Engineer", "Gurgaon", "Vector search + hybrid retrieval"),
        ("2", "Senior ML Engineer", "Noida", "Google + Flipkart career, Weaviate expert"),
        ("3", "Senior NLP Engineer", "Chennai", "Strong IR background, active on platform"),
        ("4", "Search Engineer", "Gurgaon", "7.6 yrs, ranking systems specialist"),
        ("5", "Senior ML Engineer", "Pune", "LinkedIn experience, LangChain 96.5 score"),
        ("6", "Applied ML Engineer", "Hyderabad", "Pinecone 82.0, Weaviate 85.7 assessments"),
        ("7", "Senior Data Scientist", "Delhi", "RAG systems, strong behavioral signals"),
        ("8", "Staff ML Engineer", "Kochi", "Pinecone + BM25 expert, high GitHub score"),
        ("9", "Data Scientist", "Delhi", "Semantic search focus, fast response rate"),
        ("10", "Lead AI Engineer", "Jaipur", "Razorpay, 30-day notice, IR assessment 64.8")
    ]
    
    for row_idx, data in enumerate(top_10):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_GRAY
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Segoe UI'
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            if col_idx == 0:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True
                p.font.color.rgb = CYAN

    # Slide 10: Why This Works
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide10)
    add_title(slide10, "Designed for NDCG@10 — Top 10 Is Everything")
    
    why_box = slide10.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf_why = why_box.text_frame
    tf_why.word_wrap = True
    tf_why.margin_left = tf_why.margin_top = tf_why.margin_bottom = tf_why.margin_right = 0
    
    reasons = [
        ("• Behavioral Scaling Secures Reachability", "Evaluating responsiveness and notice periods means that our top-ranked candidates are highly likely to engage, interview, and accept a fast-turnaround offer."),
        ("• Honeypot Safeguards Protect Integrity", "Filtering out resume spammers, timeline inflations, and impossible profiles ensures the shortlist contains only valid, professional, and reliable candidates."),
        ("• Location Weight Preserves Local Intent", "A 15% location penalty on non-India candidates who reject relocation aligns the candidate list with the realities of hiring for India-based offices.")
    ]
    
    for title, desc in reasons:
        p_t = tf_why.add_paragraph() if tf_why.text else tf_why.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Segoe UI'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN
        p_t.space_after = Pt(5)
        
        p_d = tf_why.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Segoe UI'
        p_d.font.size = Pt(18)
        p_d.font.color.rgb = WHITE
        p_d.space_after = Pt(25)

    # Slide 11: Tech Stack
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide11)
    add_title(slide11, "Lightweight, Fast, Reproducible")
    
    # Left Col: What we used
    left_t = slide11.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5.0))
    tf_tl = left_t.text_frame
    tf_tl.word_wrap = True
    tf_tl.margin_left = tf_tl.margin_top = tf_tl.margin_bottom = tf_tl.margin_right = 0
    
    p_uh = tf_tl.paragraphs[0]
    p_uh.text = "What We Used"
    p_uh.font.name = 'Segoe UI'
    p_uh.font.size = Pt(24)
    p_uh.font.bold = True
    p_uh.font.color.rgb = CYAN
    p_uh.space_after = Pt(20)
    
    used = [
        "Python 3.10+ Standard Library",
        "python-dateutil (flexible date parsing)",
        "tqdm (optimized execution feedback)",
        "multiprocessing (native multicore scaling)"
    ]
    for item in used:
        p_u = tf_tl.add_paragraph()
        p_u.text = f"▪ {item}"
        p_u.font.name = 'Segoe UI'
        p_u.font.size = Pt(18)
        p_u.font.color.rgb = WHITE
        p_u.space_after = Pt(15)

    # Right Col: What we avoided
    right_t = slide11.shapes.add_textbox(Inches(6.83), Inches(1.8), Inches(5.75), Inches(5.0))
    tf_tr = right_t.text_frame
    tf_tr.word_wrap = True
    tf_tr.margin_left = tf_tr.margin_top = tf_tr.margin_bottom = tf_tr.margin_right = 0
    
    p_ah = tf_tr.paragraphs[0]
    p_ah.text = "What We Deliberately Avoided"
    p_ah.font.name = 'Segoe UI'
    p_ah.font.size = Pt(24)
    p_ah.font.bold = True
    p_ah.font.color.rgb = CYAN
    p_ah.space_after = Pt(20)
    
    avoided = [
        "No external API dependencies or networks",
        "No expensive GPU prerequisites",
        "No slow LLM-based per-candidate inference",
        "Runs the full 100K candidates in < 3 minutes",
        "Peak memory footprint under 4GB RAM"
    ]
    for item in avoided:
        p_a = tf_tr.add_paragraph()
        p_a.text = f"▪ {item}"
        p_a.font.name = 'Segoe UI'
        p_a.font.size = Pt(18)
        p_a.font.color.rgb = WHITE
        p_a.space_after = Pt(15)

    # Slide 12: Team IOO
    slide12 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide12)
    add_title(slide12, "Team IOO")
    
    names_box = slide12.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf_names = names_box.text_frame
    tf_names.word_wrap = True
    tf_names.margin_left = tf_names.margin_top = tf_names.margin_bottom = tf_names.margin_right = 0
    
    p_n1 = tf_names.paragraphs[0]
    p_n1.alignment = PP_ALIGN.CENTER
    p_n1.text = "Iqra · Omar · Ojas"
    p_n1.font.name = 'Segoe UI'
    p_n1.font.size = Pt(54)
    p_n1.font.bold = True
    p_n1.font.color.rgb = WHITE
    p_n1.space_after = Pt(20)
    
    p_n2 = tf_names.add_paragraph()
    p_n2.alignment = PP_ALIGN.CENTER
    p_n2.text = "Built for India Runs at Active Skills — Data & AI Challenge 2026"
    p_n2.font.name = 'Segoe UI'
    p_n2.font.size = Pt(20)
    p_n2.font.color.rgb = CYAN
    p_n2.space_after = Pt(20)
    
    p_n3 = tf_names.add_paragraph()
    p_n3.alignment = PP_ALIGN.CENTER
    p_n3.text = "Repo: github.com/active-skills-ioo/AI_Ranking_System"
    p_n3.font.name = 'Segoe UI'
    p_n3.font.size = Pt(16)
    p_n3.font.color.rgb = LIGHT_GRAY

    # Save
    out_dir = r"C:\Projects\[PUB] India_runs_data_and_ai_challenge"
    out_path = os.path.join(out_dir, "IOO_AI_Ranking_System.pptx")
    prs.save(out_path)
    print(f"Presentation saved successfully at: {out_path}")

if __name__ == "__main__":
    create_presentation()
